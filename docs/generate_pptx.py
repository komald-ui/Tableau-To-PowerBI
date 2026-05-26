#!/usr/bin/env python3
"""Generate a Tableau-to-Power BI Migration Tool presentation (PPTX).

Produces a professional slide deck with architecture diagrams, pipeline
flow, feature highlights, migration stats, and lineage visualization.

Usage:
    python docs/generate_pptx.py [--output PATH]
"""

import argparse
import os
import sys

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Image paths ───────────────────────────────────────────────────────
_SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
IMG_MIGRATION_RESULTS = os.path.join(_SCRIPT_DIR, "images", "migration_results.png")
IMG_SHARE_ASSESSMENT = os.path.join(_SCRIPT_DIR, "images", "share_assessment.png")

# ── Color palette (PBI / Fluent) ──────────────────────────────────────
PBI_BLUE = RGBColor(0x00, 0x78, 0xD4)
PBI_DARK = RGBColor(0x00, 0x45, 0x78)
PBI_YELLOW = RGBColor(0xF2, 0xC8, 0x11)
TABLEAU_ORANGE = RGBColor(0xE9, 0x76, 0x27)
SUCCESS = RGBColor(0x10, 0x7C, 0x10)
PURPLE = RGBColor(0x87, 0x64, 0xB8)
TEAL = RGBColor(0x03, 0x83, 0x87)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x24, 0x29, 0x2E)
LIGHT_GRAY = RGBColor(0xF3, 0xF2, 0xF1)
GRAY = RGBColor(0x60, 0x5E, 0x5C)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def _add_gradient_bg(slide, color_top, color_bottom):
    """Add a gradient background fill to a slide."""
    bg = slide.background
    fill = bg.fill
    fill.gradient()
    fill.gradient_stops[0].color.rgb = color_top
    fill.gradient_stops[0].position = 0.0
    fill.gradient_stops[1].color.rgb = color_bottom
    fill.gradient_stops[1].position = 1.0


def _add_box(slide, left, top, width, height, fill_color, text, font_size=14,
             font_color=WHITE, bold=False, border_color=None):
    """Add a rounded rectangle with centered text."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(2)
    else:
        shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = font_color
    run.font.bold = bold
    tf.paragraphs[0].space_before = Pt(0)
    tf.paragraphs[0].space_after = Pt(0)
    # Vertical center
    tf.word_wrap = True
    try:
        tf.auto_size = None
        shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        shape.text_frame_anchor = MSO_ANCHOR.MIDDLE
    except Exception:
        pass
    return shape


def _add_arrow(slide, start_left, start_top, end_left, end_top, color=GRAY):
    """Add a connector arrow between two points."""
    connector = slide.shapes.add_connector(
        1,  # straight connector
        start_left, start_top, end_left, end_top)
    connector.line.color.rgb = color
    connector.line.width = Pt(2)
    # Arrow head
    connector.end_x = end_left
    connector.end_y = end_top
    return connector


def _add_text(slide, left, top, width, height, text, font_size=18,
              font_color=DARK, bold=False, alignment=PP_ALIGN.LEFT):
    """Add a text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = font_color
    run.font.bold = bold
    return txBox


def _add_bullet_list(slide, left, top, width, height, items, font_size=16,
                     font_color=DARK):
    """Add a bulleted text list."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4)
        p.space_after = Pt(4)
        run = p.add_run()
        run.text = f"  •  {item}"
        run.font.size = Pt(font_size)
        run.font.color.rgb = font_color
    return txBox


# ══════════════════════════════════════════════════════════════════════
#  SLIDE BUILDERS
# ══════════════════════════════════════════════════════════════════════

def slide_title(prs):
    """Slide 1: Title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _add_gradient_bg(slide, PBI_DARK, PBI_BLUE)

    _add_text(slide, Inches(1.5), Inches(1.5), Inches(10), Inches(1.2),
              "Tableau → Power BI", 48, WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    _add_text(slide, Inches(1.5), Inches(2.7), Inches(10), Inches(0.8),
              "Automated Migration Tool", 36, PBI_YELLOW, bold=True, alignment=PP_ALIGN.CENTER)
    _add_text(slide, Inches(1.5), Inches(4.0), Inches(10), Inches(0.6),
              "Migrate .twb/.twbx workbooks to .pbip projects — fully automated, zero manual rework",
              18, WHITE, alignment=PP_ALIGN.CENTER)

    # Stats bar
    stats = [
        ("8,668", "Tests Passed"),
        ("133+", "DAX Conversions"),
        ("190", "Visual Types"),
        ("49", "Data Connectors"),
        ("100%", "Fidelity (27/27 WB)"),
    ]
    bar_left = Inches(1.0)
    bar_top = Inches(5.5)
    box_w = Inches(2.0)
    box_h = Inches(1.2)
    gap = Inches(0.27)
    for i, (val, label) in enumerate(stats):
        x = bar_left + i * (box_w + gap)
        shape = _add_box(slide, x, bar_top, box_w, box_h, PBI_DARK,
                         f"{val}\n{label}", 16, WHITE, bold=True)
        shape.fill.fore_color.rgb = RGBColor(0x00, 0x35, 0x60)


def slide_executive_summary(prs):
    """Slide 2: Executive Summary — Functionality, Optimization & ROI."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

    _add_text(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.6),
              "Executive Summary — End-to-End Value", 32, PBI_DARK, bold=True)

    # ── Left column: Functionality layers diagram ──
    _add_text(slide, Inches(0.3), Inches(0.9), Inches(5.5), Inches(0.4),
              "Full-Stack Migration Pipeline", 18, PBI_DARK, bold=True)

    layers = [
        ("📥 INPUT", "TWB / TWBX / TDS / Prep / Hyper / Server API", TABLEAU_ORANGE),
        ("🔍 EXTRACT", "23 object types → 23 JSON intermediate files", PBI_BLUE),
        ("🧮 CONVERT", "133+ DAX formulas • 43 M transforms • 49 connectors", PURPLE),
        ("⚙️ GENERATE", "TMDL semantic model • PBIR v4.0 report • themes", PBI_BLUE),
        ("🧠 OPTIMIZE", "DAX AST rewriter • Time Intelligence inject • auto-fix", TEAL),
        ("✅ VALIDATE", "QA suite • governance • lineage • comparison report", SUCCESS),
        ("🚀 DEPLOY", "PBI Service • Fabric • bundle • multi-tenant • gateway", PBI_DARK),
    ]

    layer_w = Inches(5.5)
    layer_h = Inches(0.62)
    layer_x = Inches(0.3)
    layer_y_start = Inches(1.35)
    layer_gap = Inches(0.08)

    for i, (label, desc, color) in enumerate(layers):
        y = layer_y_start + i * (layer_h + layer_gap)
        _add_box(slide, layer_x, y, layer_w, layer_h, color,
                 f"{label}  —  {desc}", 11, WHITE, bold=False)
        # Down arrow between layers
        if i < len(layers) - 1:
            ay = y + layer_h
            _add_text(slide, Inches(2.7), ay - Inches(0.05), Inches(0.5), Inches(0.2),
                      "▼", 10, GRAY, alignment=PP_ALIGN.CENTER)

    # ── Right column top: Optimization at every level ──
    _add_text(slide, Inches(6.2), Inches(0.9), Inches(6.5), Inches(0.4),
              "Optimization at Every Level", 18, PBI_DARK, bold=True)

    optimizations = [
        ("Data Layer", "Smart M queries • incremental refresh\nconnection parameterization", PBI_BLUE),
        ("Semantic Layer", "DAX optimizer (IF→SWITCH • COALESCE)\nTime Intelligence auto-inject\nSortByColumn • displayFolder • isHidden", PURPLE),
        ("Visual Layer", "190 visual mappings • conditional format\ntheme extraction • responsive layout\nreference lines • tooltips", TEAL),
        ("Security Layer", "RLS roles • USERPRINCIPALNAME()\nZIP slip / XXE defense\ncredential redaction", RGBColor(0xA4, 0x26, 0x2C)),
    ]

    opt_w = Inches(3.3)
    opt_h = Inches(1.15)
    opt_gap_x = Inches(0.15)
    opt_gap_y = Inches(0.15)

    for i, (title, desc, color) in enumerate(optimizations):
        col = i % 2
        row = i // 2
        x = Inches(6.2) + col * (opt_w + opt_gap_x)
        y = Inches(1.35) + row * (opt_h + opt_gap_y)
        _add_box(slide, x, y, opt_w, opt_h, color,
                 f"{title}\n{desc}", 10, WHITE, bold=False)

    # ── Right column bottom: ROI metrics ──
    _add_text(slide, Inches(6.2), Inches(4.0), Inches(6.5), Inches(0.4),
              "Migration ROI & Business Impact", 18, PBI_DARK, bold=True)

    roi_cards = [
        ("⏱️", "95%+", "Time Saved", "vs. manual rework"),
        ("💰", "~$0", "Tooling Cost", "Python stdlib only"),
        ("🎯", "100%", "Fidelity", "27/27 workbooks"),
        ("🔒", "0", "Dependencies", "No 3rd-party risk"),
    ]

    roi_w = Inches(1.6)
    roi_h = Inches(1.4)
    roi_gap = Inches(0.1)
    roi_x_start = Inches(6.2)
    roi_y = Inches(4.5)

    for i, (icon, val, label, sub) in enumerate(roi_cards):
        x = roi_x_start + i * (roi_w + roi_gap)
        shape = _add_box(slide, x, roi_y, roi_w, roi_h, LIGHT_GRAY,
                         f"{icon}\n{val}\n{label}\n{sub}", 11, DARK, bold=False,
                         border_color=PBI_BLUE)

    # ROI summary line
    _add_text(slide, Inches(0.3), Inches(6.6), Inches(12.5), Inches(0.5),
              "Zero external dependencies • Zero manual rework on standard workbooks • "
              "Enterprise-grade security (OWASP Top 10) • Full audit trail & governance",
              13, GRAY, alignment=PP_ALIGN.CENTER)


def slide_pipeline(prs):
    """Slide 2: 2-Step Migration Pipeline diagram."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

    _add_text(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
              "Migration Pipeline — 2-Step Architecture", 32, PBI_DARK, bold=True)

    # Step boxes
    box_w = Inches(2.5)
    box_h = Inches(1.6)
    y = Inches(2.0)

    # Source
    _add_box(slide, Inches(0.5), y, box_w, box_h, TABLEAU_ORANGE,
             "📄 Tableau Source\n.twb / .twbx / .tds\nPrep flows (.tfl)", 14, WHITE, bold=True)

    # Arrow 1
    _add_text(slide, Inches(3.1), Inches(2.5), Inches(0.8), Inches(0.5),
              "→", 36, GRAY, bold=True, alignment=PP_ALIGN.CENTER)

    # Extract
    _add_box(slide, Inches(3.8), y, box_w, box_h, PBI_BLUE,
             "🔍 EXTRACT\nParse XML → 23 JSON\nfiles (worksheets,\ncalcs, datasources…)", 13, WHITE, bold=True)

    # Arrow 2
    _add_text(slide, Inches(6.4), Inches(2.5), Inches(0.8), Inches(0.5),
              "→", 36, GRAY, bold=True, alignment=PP_ALIGN.CENTER)

    # Generate
    _add_box(slide, Inches(7.1), y, box_w, box_h, PBI_BLUE,
             "⚙️ GENERATE\nTMDL Semantic Model\nPBIR v4.0 Report\nPower Query M", 13, WHITE, bold=True)

    # Arrow 3
    _add_text(slide, Inches(9.7), Inches(2.5), Inches(0.8), Inches(0.5),
              "→", 36, GRAY, bold=True, alignment=PP_ALIGN.CENTER)

    # Output
    _add_box(slide, Inches(10.4), y, box_w, box_h, SUCCESS,
             "📊 Power BI\n.pbip project\nOpen in PBI Desktop\n(March 2025+)", 14, WHITE, bold=True)

    # Optional deploy path
    _add_text(slide, Inches(10.4), Inches(3.8), Inches(2.5), Inches(0.5),
              "↓  optional", 12, GRAY, alignment=PP_ALIGN.CENTER)
    _add_box(slide, Inches(10.4), Inches(4.3), box_w, Inches(1.0), PBI_YELLOW,
             "🚀 Deploy to PBI Service / Fabric", 13, DARK, bold=True)

    # Extracted objects list
    _add_text(slide, Inches(0.5), Inches(4.3), Inches(9), Inches(0.5),
              "23 Extracted Object Types:", 16, PBI_DARK, bold=True)
    objects = [
        "worksheets • dashboards • datasources • calculations • parameters • filters",
        "stories • actions • sets • groups • bins • hierarchies • sort orders",
        "aliases • custom SQL • user filters (→ RLS) • hyper file data",
    ]
    for i, line in enumerate(objects):
        _add_text(slide, Inches(0.5), Inches(4.9 + i * 0.35), Inches(9), Inches(0.35),
                  line, 13, GRAY)

    # Fabric path
    _add_box(slide, Inches(7.1), Inches(4.3), box_w, Inches(1.6), RGBColor(0x00, 0x78, 0xD4),
             "🏭 Fabric Mode\nLakehouse + Dataflow\nNotebook + DirectLake\n+ Pipeline", 12, WHITE, bold=True)
    _add_text(slide, Inches(7.1), Inches(6.0), Inches(2.5), Inches(0.4),
              "--output-format fabric", 11, GRAY, alignment=PP_ALIGN.CENTER)


def slide_features(prs):
    """Slide 3: Key Features grid."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

    _add_text(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
              "Key Features", 32, PBI_DARK, bold=True)

    features = [
        ("🔄 Complete Extraction", "23 object types from .twb/.twbx\nPrep flows, Hyper files, Server API", PBI_BLUE),
        ("🧮 133+ DAX Conversions", "LOD, table calcs, RANK, WINDOW\ncross-table RELATED/LOOKUPVALUE", PURPLE),
        ("📊 190 Visual Types", "Bar, line, map, treemap, Sankey\ncombo, gauge, word cloud, KPI", TEAL),
        ("🔌 49 Data Connectors", "SQL Server, Snowflake, BigQuery\nDatabricks, Oracle, Excel, CSV…", PBI_BLUE),
        ("🧠 Smart Semantic Model", "Calendar, hierarchies, calc groups\nfield params, RLS, perspectives\nM identifier auto-quoting", PURPLE),
        ("⚡ DAX Optimizer", "IF→SWITCH, COALESCE, constant\nfolding, Time Intelligence inject", TEAL),
        ("🔗 Shared Semantic Model", "Multi-WB merge, fingerprint match\nJaccard scoring, thin reports", PBI_BLUE),
        ("🔍 QA Suite & Auto-Fix", "17 auto-fix patterns, governance\nlineage map, comparison reports", PURPLE),
        ("🏭 Fabric-Native Output", "Lakehouse, Dataflow Gen2, PySpark\nDirectLake, Data Pipeline", TEAL),
        ("🚀 Deploy Anywhere", "PBI Service, Fabric, gateway\nbundle deploy, multi-tenant", PBI_BLUE),
        ("📈 Assessment & Strategy", "9-category readiness scoring\nImport/DQ/Composite advisor", PURPLE),
        ("🔗 Lineage Map", "Object traceability: source→target\nflow diagrams, HTML dashboard", TEAL),
    ]

    cols = 4
    rows = 3
    box_w = Inches(2.9)
    box_h = Inches(1.5)
    x_start = Inches(0.5)
    y_start = Inches(1.4)
    gap_x = Inches(0.25)
    gap_y = Inches(0.25)

    for i, (title, desc, color) in enumerate(features):
        row = i // cols
        col = i % cols
        x = x_start + col * (box_w + gap_x)
        y = y_start + row * (box_h + gap_y)
        _add_box(slide, x, y, box_w, box_h, color,
                 f"{title}\n\n{desc}", 12, WHITE, bold=False)


def slide_dax(prs):
    """Slide 4: DAX Conversion highlights."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

    _add_text(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
              "133+ DAX Conversions — Highlights", 32, PBI_DARK, bold=True)

    conversions = [
        ("LOD Fixed", "{FIXED [customer] : SUM([qty])}", "CALCULATE(SUM([qty]),\n  ALLEXCEPT('T','T'[customer]))"),
        ("LOD Exclude", "{EXCLUDE [channel] : SUM([rev])}", "CALCULATE(SUM([rev]),\n  REMOVEFILTERS('T'[channel]))"),
        ("Iterator", "SUM(IF [status]=\"OK\" THEN [qty])", "SUMX('Tbl',\n  IF([status]=\"OK\", [qty], 0))"),
        ("Table Calc", "RUNNING_SUM(SUM([sales]))", "CALCULATE(SUM([sales]),\n  FILTER(ALL('T'),…))"),
        ("Rank", "RANK(SUM([revenue]))", "RANKX(ALL('T'),\n  SUM('T'[revenue]))"),
        ("Cross-Table", "[col] from other table", "RELATED('Other'[col]) or\n  LOOKUPVALUE(…)"),
        ("Security", "USERNAME()", "USERPRINCIPALNAME()"),
        ("Null", "ISNULL / ZN / IFNULL", "ISBLANK / IF(ISBLANK,…)"),
    ]

    headers = ["Category", "Tableau", "DAX"]
    col_widths = [Inches(1.8), Inches(4.5), Inches(5.5)]
    x_start = Inches(0.5)
    y_start = Inches(1.3)
    row_h = Inches(0.7)

    # Header row
    x = x_start
    for j, (hdr, w) in enumerate(zip(headers, col_widths)):
        _add_box(slide, x, y_start, w, Inches(0.5), PBI_DARK, hdr, 14, WHITE, bold=True)
        x += w

    # Data rows
    for i, (cat, tableau, dax) in enumerate(conversions):
        y = y_start + Inches(0.55) + i * row_h
        bg = LIGHT_GRAY if i % 2 == 0 else WHITE
        x = x_start
        _add_box(slide, x, y, col_widths[0], row_h, bg, cat, 12, PBI_BLUE, bold=True, border_color=LIGHT_GRAY)
        x += col_widths[0]
        _add_box(slide, x, y, col_widths[1], row_h, bg, tableau, 11, TABLEAU_ORANGE, bold=False, border_color=LIGHT_GRAY)
        x += col_widths[1]
        _add_box(slide, x, y, col_widths[2], row_h, bg, dax, 11, SUCCESS, bold=False, border_color=LIGHT_GRAY)


def slide_lineage(prs):
    """Slide 5: Lineage Map & Traceability."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

    _add_text(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
              "Lineage Map — Source-to-Target Traceability", 32, PBI_DARK, bold=True)

    _add_text(slide, Inches(0.5), Inches(1.1), Inches(12), Inches(0.6),
              "Every migration produces lineage_map.json — visualized in the HTML dashboard",
              16, GRAY)

    # Flow diagram
    flow_boxes = [
        ("Tableau\nDatasource.Table", TABLEAU_ORANGE),
        ("Calculations\n(formulas)", PURPLE),
        ("PBI Semantic Model\n(tables, measures,\nrelationships)", PBI_BLUE),
        ("Report Pages\n(visuals, filters)", SUCCESS),
    ]
    box_w = Inches(2.5)
    box_h = Inches(1.3)
    y = Inches(2.0)
    gap = Inches(0.5)
    arrow_w = Inches(0.6)

    for i, (label, color) in enumerate(flow_boxes):
        x = Inches(0.5) + i * (box_w + arrow_w)
        _add_box(slide, x, y, box_w, box_h, color, label, 14, WHITE, bold=True)
        if i < len(flow_boxes) - 1:
            ax = Inches(0.5) + (i + 1) * box_w + i * arrow_w + Inches(0.05)
            _add_text(slide, ax, Inches(2.3), arrow_w, Inches(0.5),
                      "→", 32, PBI_BLUE, bold=True, alignment=PP_ALIGN.CENTER)

    # Lineage categories
    categories = [
        ("📦 Tables", "Tableau datasource.table\n→ PBI table name", PBI_BLUE),
        ("🧮 Calculations", "Tableau calc formula\n→ PBI measure/column + type", PURPLE),
        ("🔗 Relationships", "From Table[Col]\n→ To Table[Col] + cardinality", TEAL),
        ("📊 Worksheets", "Tableau worksheet name\n→ PBI report page", SUCCESS),
    ]
    cat_w = Inches(2.8)
    cat_h = Inches(1.5)
    cat_y = Inches(4.0)
    for i, (title, desc, color) in enumerate(categories):
        x = Inches(0.5) + i * (cat_w + Inches(0.27))
        _add_box(slide, x, cat_y, cat_w, cat_h, color,
                 f"{title}\n\n{desc}", 13, WHITE, bold=False)

    _add_text(slide, Inches(0.5), Inches(5.8), Inches(12), Inches(0.5),
              "HTML Dashboard: searchable/sortable tables per category • flow diagrams • stat cards • multi-workbook aggregation",
              14, GRAY)


def slide_shared_model(prs):
    """Slide 6: Shared Semantic Model & Merge."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

    _add_text(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
              "Shared Semantic Model — Multi-Workbook Merge", 32, PBI_DARK, bold=True)

    # Input workbooks
    wb_y = Inches(1.5)
    for i, name in enumerate(["Workbook A", "Workbook B", "Workbook C"]):
        _add_box(slide, Inches(0.5), wb_y + i * Inches(1.3), Inches(2.2), Inches(1.0),
                 TABLEAU_ORANGE, f"📄 {name}", 14, WHITE, bold=True)
        _add_text(slide, Inches(2.8), wb_y + i * Inches(1.3) + Inches(0.2),
                  Inches(0.8), Inches(0.5), "→", 28, GRAY, bold=True, alignment=PP_ALIGN.CENTER)

    # Merge engine
    _add_box(slide, Inches(3.5), Inches(2.0), Inches(2.8), Inches(2.0), PBI_BLUE,
             "🔗 MERGE\n\nFingerprint matching\nJaccard scoring\nConflict resolution", 13, WHITE, bold=True)

    # Output
    _add_text(slide, Inches(6.4), Inches(2.7), Inches(0.8), Inches(0.5),
              "→", 28, GRAY, bold=True, alignment=PP_ALIGN.CENTER)

    _add_box(slide, Inches(7.2), Inches(1.5), Inches(2.8), Inches(1.2), PBI_DARK,
             "📦 Shared\nSemantic Model\n(ONE model)", 14, WHITE, bold=True)

    for i, name in enumerate(["Report A (thin)", "Report B (thin)", "Report C (thin)"]):
        _add_box(slide, Inches(10.2), Inches(1.3) + i * Inches(1.1), Inches(2.5), Inches(0.8),
                 PBI_YELLOW, f"📊 {name}", 12, DARK, bold=True)
        # byPath arrow
        _add_text(slide, Inches(10.0), Inches(1.3) + i * Inches(1.1) + Inches(0.15),
                  Inches(0.4), Inches(0.4), "←", 18, GRAY, bold=True, alignment=PP_ALIGN.CENTER)

    _add_text(slide, Inches(10.2), Inches(4.5), Inches(2.5), Inches(0.4),
              "byPath → SharedModel", 11, GRAY, alignment=PP_ALIGN.CENTER)

    # Features list
    features = [
        "SHA-256 fingerprint-based table matching across workbooks",
        "Jaccard column overlap scoring (0–100 merge score)",
        "Measure, column, relationship deduplication & conflict resolution",
        "Global assessment: pairwise merge scoring + BFS clustering",
        "Bundle deployment: shared model + thin reports → Fabric workspace",
    ]
    _add_bullet_list(slide, Inches(0.5), Inches(5.0), Inches(12), Inches(2.0),
                     features, 14, DARK)

    # Share assessment screenshot
    if os.path.isfile(IMG_SHARE_ASSESSMENT):
        slide.shapes.add_picture(IMG_SHARE_ASSESSMENT,
                                 Inches(3.5), Inches(4.3), Inches(9.0))


def slide_migration_results(prs):
    """Slide 7: Migration Results (27/27 workbooks)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

    _add_text(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
              "Migration Results — 27/27 Workbooks at 100% Fidelity", 32, PBI_DARK, bold=True)

    # Stats cards
    cards = [
        ("27/27", "Workbooks\nSucceeded", SUCCESS),
        ("0", "Failed", RGBColor(0xA4, 0x26, 0x2C)),
        ("100%", "Average\nFidelity", PBI_BLUE),
        ("7,099", "Tests\nPassed", PURPLE),
        ("0", "Test\nFailures", SUCCESS),
    ]
    card_w = Inches(2.2)
    card_h = Inches(1.5)
    x_start = Inches(0.5)
    gap = Inches(0.3)
    for i, (val, label, color) in enumerate(cards):
        x = x_start + i * (card_w + gap)
        _add_box(slide, x, Inches(1.3), card_w, card_h, color,
                 f"{val}\n{label}", 18, WHITE, bold=True)

    # Tables: Samples + Real World
    _add_text(slide, Inches(0.5), Inches(3.2), Inches(5.5), Inches(0.5),
              "Tableau Samples (10 workbooks)", 18, PBI_DARK, bold=True)
    samples = [
        "Customer_Analysis", "Financial_Report", "Global_Superstores",
        "HR_Dashboard", "Marketing_Analytics", "Sales_Dashboard",
        "Sales_Performance", "Supply_Chain", "Tech_Support", "Ventes_France",
    ]
    for i, name in enumerate(samples):
        row = i // 5
        col = i % 5
        _add_box(slide, Inches(0.5) + col * Inches(1.2), Inches(3.8) + row * Inches(0.5),
                 Inches(1.1), Inches(0.4), LIGHT_GRAY, name[:16], 9, SUCCESS, bold=True, border_color=LIGHT_GRAY)

    _add_text(slide, Inches(6.5), Inches(3.2), Inches(6.5), Inches(0.5),
              "Real-World Workbooks (17 workbooks)", 18, PBI_DARK, bold=True)
    real_world = [
        "Cache", "RESTAPISample", "SampleDS", "SampleWB", "TABLEAU_10_TWB",
        "TABLEAU_10_TWBX", "datasource_test", "ephemeral_field",
        "feedback_dashboard", "filtering", "global_superstores_db",
        "multiple_connections", "nba_player_stats", "sample-superstore",
        "shapes_test", "superstore_sales", "vishnu_dashboard",
    ]
    for i, name in enumerate(real_world):
        row = i // 6
        col = i % 6
        _add_box(slide, Inches(6.5) + col * Inches(1.15), Inches(3.8) + row * Inches(0.5),
                 Inches(1.08), Inches(0.4), LIGHT_GRAY, name[:16], 8, SUCCESS, bold=True, border_color=LIGHT_GRAY)

    _add_text(slide, Inches(0.5), Inches(5.5), Inches(12), Inches(0.5),
              "All 27 workbooks migrated at 100% fidelity — zero errors, zero manual intervention required.",
              16, SUCCESS, bold=True, alignment=PP_ALIGN.CENTER)

    # Migration results screenshot
    if os.path.isfile(IMG_MIGRATION_RESULTS):
        slide.shapes.add_picture(IMG_MIGRATION_RESULTS,
                                 Inches(6.5), Inches(5.3), Inches(6.5))


def slide_qa_automation(prs):
    """Slide 8: QA & Post-Migration Automation."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

    _add_text(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
              "Post-Migration Automation & QA", 32, PBI_DARK, bold=True)

    # QA Pipeline flow
    qa_steps = [
        ("Validate\n(.pbip JSON/TMDL)", PBI_BLUE),
        ("Auto-Fix\n(17 DAX patterns)", PURPLE),
        ("Governance\n(naming, PII, labels)", TEAL),
        ("Compare\n(Tableau vs PBI)", PBI_BLUE),
        ("qa_report.json\n(pass/fail/score)", SUCCESS),
    ]
    box_w = Inches(2.2)
    box_h = Inches(1.2)
    y = Inches(1.5)
    for i, (label, color) in enumerate(qa_steps):
        x = Inches(0.3) + i * (box_w + Inches(0.25))
        _add_box(slide, x, y, box_w, box_h, color, label, 13, WHITE, bold=True)
        if i < len(qa_steps) - 1:
            ax = Inches(0.3) + (i + 1) * box_w + i * Inches(0.25) + Inches(0.02)
            _add_text(slide, ax, Inches(1.7), Inches(0.25), Inches(0.5),
                      "→", 22, GRAY, bold=True, alignment=PP_ALIGN.CENTER)

    _add_text(slide, Inches(0.5), Inches(2.9), Inches(5), Inches(0.4),
              "python migrate.py workbook.twbx --qa", 14, PBI_BLUE, bold=True)

    # Automation features
    auto_features = [
        ("🔧 Validator Auto-Fix", "17 Tableau→DAX leak patterns:\nISNULL→ISBLANK, ZN→IF(ISBLANK),\nELSEIF→IF, ==→=, or→||, and→&&"),
        ("📜 RLS PowerShell", "Auto-generate .ps1 scripts for\nAzure AD RLS role assignment\nvia Power BI REST API"),
        ("🔑 Credential Template", "JSON placeholder file for each\ndatasource connection — ready\nfor DevOps pipeline injection"),
        ("🔗 Lineage Map", "lineage_map.json tracks every\nobject: tables, calcs, rels,\nworksheets → HTML dashboard"),
    ]

    for i, (title, desc) in enumerate(auto_features):
        col = i % 4
        x = Inches(0.5) + col * (Inches(3.0) + Inches(0.2))
        y = Inches(3.6)
        _add_box(slide, x, y, Inches(3.0), Inches(2.0), LIGHT_GRAY,
                 f"{title}\n\n{desc}", 12, DARK, bold=False, border_color=PBI_BLUE)


def slide_testing(prs):
    """Slide 9: Testing & CI/CD."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

    _add_text(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
              "Testing & CI/CD Pipeline", 32, PBI_DARK, bold=True)

    # CI steps
    ci_steps = [
        ("🔍 Lint\nflake8 + ruff", PURPLE),
        ("🧪 Test\n7,099 tests\nPy 3.12–3.14", SUCCESS),
        ("✅ Validate\nStrict .twbx\nmigrations", PBI_BLUE),
        ("📦 Staging\nFabric deploy", PBI_YELLOW),
        ("🚀 Production\nManual approval", RGBColor(0xEF, 0x44, 0x44)),
    ]
    box_w = Inches(2.2)
    box_h = Inches(1.4)
    y = Inches(1.3)
    for i, (label, color) in enumerate(ci_steps):
        x = Inches(0.3) + i * (box_w + Inches(0.25))
        font_c = DARK if color == PBI_YELLOW else WHITE
        _add_box(slide, x, y, box_w, box_h, color, label, 14, font_c, bold=True)
        if i < len(ci_steps) - 1:
            ax = Inches(0.3) + (i + 1) * box_w + i * Inches(0.25) + Inches(0.02)
            _add_text(slide, ax, Inches(1.6), Inches(0.25), Inches(0.5),
                      "→", 22, GRAY, bold=True, alignment=PP_ALIGN.CENTER)

    # Test breakdown
    test_data = [
        ("DAX Coverage", "168 tests", "Edge cases across all DAX categories"),
        ("Generation", "145 tests", "TMDL/PBIR generation edge cases"),
        ("M Query", "102 tests", "Power Query M, 40+ transforms"),
        ("Semantic Model", "92 tests", "Calendar, TMDL, hierarchies"),
        ("Visual Generator", "65 tests", "118+ types, sync, buttons"),
        ("Real-World E2E", "63 tests", "End-to-end sample migrations"),
        ("Automation", "64 tests", "Auto-fix, lineage, QA, governance"),
        ("+ 134 more files", "6,400 tests", "Sprint, coverage, wizard, telemetry…"),
    ]

    _add_text(slide, Inches(0.5), Inches(3.0), Inches(4), Inches(0.5),
              "Test Suite Breakdown (141 files):", 18, PBI_DARK, bold=True)

    headers = ["Category", "Count", "Coverage"]
    col_ws = [Inches(2.5), Inches(1.5), Inches(4.5)]
    x_start = Inches(0.5)
    y_start = Inches(3.6)
    rh = Inches(0.42)

    x = x_start
    for j, (hdr, w) in enumerate(zip(headers, col_ws)):
        _add_box(slide, x, y_start, w, Inches(0.38), PBI_DARK, hdr, 12, WHITE, bold=True)
        x += w

    for i, (cat, count, desc) in enumerate(test_data):
        y = y_start + Inches(0.4) + i * rh
        bg = LIGHT_GRAY if i % 2 == 0 else WHITE
        x = x_start
        _add_box(slide, x, y, col_ws[0], rh, bg, cat, 11, DARK, bold=True, border_color=LIGHT_GRAY)
        x += col_ws[0]
        _add_box(slide, x, y, col_ws[1], rh, bg, count, 11, PBI_BLUE, bold=True, border_color=LIGHT_GRAY)
        x += col_ws[1]
        _add_box(slide, x, y, col_ws[2], rh, bg, desc, 11, GRAY, bold=False, border_color=LIGHT_GRAY)


def slide_next_steps(prs):
    """Slide 10: Next Steps / Call to Action."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_gradient_bg(slide, PBI_DARK, PBI_BLUE)

    _add_text(slide, Inches(1.5), Inches(0.8), Inches(10), Inches(0.8),
              "Next Steps", 40, WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    steps = [
        "1.  Open the generated .pbip in Power BI Desktop (Developer Mode)",
        "2.  Configure data source credentials in Power Query Editor",
        "3.  Verify DAX measures and relationships in Model view",
        "4.  Compare visuals side-by-side with original Tableau workbook",
        "5.  Run --qa for automated quality assurance + auto-fix",
        "6.  Deploy to Power BI Service or Microsoft Fabric",
    ]
    for i, step in enumerate(steps):
        _add_text(slide, Inches(2.0), Inches(2.0) + i * Inches(0.7),
                  Inches(9), Inches(0.6), step, 20, WHITE)

    # Command box
    _add_box(slide, Inches(2.0), Inches(6.0), Inches(9), Inches(0.7),
             RGBColor(0x00, 0x35, 0x60),
             "python migrate.py your_workbook.twbx --qa --deploy WORKSPACE_ID",
             16, PBI_YELLOW, bold=True)


# ══════════════════════════════════════════════════════════════════════
#  MAIN
# ══════════════════════════════════════════════════════════════════════

def generate_presentation(output_path):
    """Generate the full PPTX presentation."""
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    slide_title(prs)
    slide_executive_summary(prs)
    slide_pipeline(prs)
    slide_features(prs)
    slide_dax(prs)
    slide_lineage(prs)
    slide_shared_model(prs)
    slide_migration_results(prs)
    slide_qa_automation(prs)
    slide_testing(prs)
    slide_next_steps(prs)

    prs.save(output_path)
    print(f"Presentation saved: {output_path}")
    print(f"  Slides: {len(prs.slides)}")
    return output_path


def main():
    parser = argparse.ArgumentParser(description="Generate migration PPTX presentation")
    parser.add_argument("--output", "-o", default="docs/Tableau_to_PowerBI_Migration.pptx",
                        help="Output PPTX path")
    args = parser.parse_args()
    generate_presentation(args.output)


if __name__ == "__main__":
    main()
